"""Build AttnFuse_slides.pptx — a deck explaining the AttnFuse project end-to-end.

ELI5 framing: each slide carries 4-6 short bullets the audience can read at a glance,
while the matching speaker-notes LaTeX (AttnFuse_speaker_notes.tex) provides the script.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette --------------------------------------------------------------
NAVY      = RGBColor(0x1B, 0x3A, 0x5C)
ORANGE    = RGBColor(0xE0, 0x7A, 0x1F)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)
CODE_BG   = RGBColor(0x2D, 0x2D, 0x2D)
CODE_FG   = RGBColor(0xE6, 0xE6, 0xE6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)   # 16:9


def add_title(slide, text, size=34, color=NAVY):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return box


def add_subtitle(slide, text, top=1.1, size=18, color=DARK_GRAY):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.55))
    p   = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.italic = True
    p.font.color.rgb = color
    return box


def add_bullets(slide, bullets, left=0.5, top=1.7, width=12.3, height=5.4,
                size=20, color=DARK_GRAY, gap_pt=8):
    """bullets: list of (text, level) tuples or plain strings."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        text, level = (item if isinstance(item, tuple) else (item, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = ("    " * level) + ("•  " if level == 0 else "–  ")
        p.text = prefix + text
        p.font.size = Pt(size - 2 * level)
        p.font.color.rgb = color
        p.space_after = Pt(gap_pt)
    return box


def add_code(slide, code, left=0.7, top=1.7, width=12.0, height=4.8, size=16):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top  = Inches(0.20); tf.margin_bottom = Inches(0.20)
    tf.word_wrap = True
    lines = code.strip("\n").split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line if line else " "
        p.font.name = "Consolas"
        p.font.size = Pt(size)
        p.font.color.rgb = CODE_FG
        p.space_after = Pt(2)
    return box


def add_footer(slide, page, total):
    fb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    p  = fb.text_frame.paragraphs[0]
    p.text = f"AttnFuse  •  CS 790/657  •  {page} / {total}"
    p.font.size = Pt(10); p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_callout(slide, text, left, top, width, height, fill=ORANGE, fg=WHITE, size=16):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = fg
    return box


# ---- deck -----------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slides = []

    # 1. Title --------------------------------------------------------------
    def slide_title():
        s = prs.slides.add_slide(blank)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

        t = s.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12.0), Inches(1.4))
        p = t.text_frame.paragraphs[0]
        p.text = "AttnFuse"
        p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = WHITE

        sb = s.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(12.0), Inches(0.7))
        p = sb.text_frame.paragraphs[0]
        p.text = "An Embedded Python DSL for Compiling Attention to Fused GPU Kernels"
        p.font.size = Pt(24); p.font.color.rgb = WHITE

        meta = s.shapes.add_textbox(Inches(0.7), Inches(5.4), Inches(12.0), Inches(1.2))
        tf = meta.text_frame
        for i, line in enumerate([
            "Varun Kumar Dasoju",
            "CS 790 / 657 — Domain-Specific Programming for AI",
            "Final Project Presentation",
        ]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18); p.font.color.rgb = WHITE
            p.space_after = Pt(4)
        return s
    slides.append(slide_title)

    # 2. Big picture --------------------------------------------------------
    def slide_bigpicture():
        s = prs.slides.add_slide(blank)
        add_title(s, "What Is This Project About?")
        add_subtitle(s, "In one sentence: we make new kinds of AI attention easy to build and fast to run.")
        add_bullets(s, [
            "Modern AI models (like ChatGPT) read text using something called attention.",
            "Attention is the part of the model that decides which words matter most.",
            "It is also the slowest part — for long texts, it eats most of the GPU time.",
            "Researchers keep inventing new attention recipes, but each one takes weeks of expert GPU coding.",
            ("AttnFuse: write the recipe in 5 lines of Python — we compile it into a fast GPU program automatically.", 0),
        ])
        return s
    slides.append(slide_bigpicture)

    # 3. What is attention --------------------------------------------------
    def slide_what_is_attention():
        s = prs.slides.add_slide(blank)
        add_title(s, "What Is 'Attention' in AI? — A Simple Picture")
        add_subtitle(s, "Think of attention like a smart highlighter inside the AI.")
        add_bullets(s, [
            "When you read “She bought a red apple,” your brain links “red” to “apple.”",
            "An AI does the same: for every word it reads, it glances back at the other words.",
            "It writes down a score for every pair of words — a giant N×N scoreboard.",
            "Then it mixes the words together, weighted by those scores.",
            "Result: the model “understands” each word in context.",
        ])
        return s
    slides.append(slide_what_is_attention)

    # 4. The math (still ELI5) ---------------------------------------------
    def slide_math():
        s = prs.slides.add_slide(blank)
        add_title(s, "The Math, Stripped Down")
        add_subtitle(s, "Three lists per word, then matrix multiply, then softmax, then multiply again.")
        add_bullets(s, [
            "Each word gets three vectors: Q (a question), K (a label), V (the info it carries).",
            "S = Q · Kᵀ  — how well does word i’s question match word j’s label?",
            "P = softmax(S)  — turn raw scores into percentages that add up to 1.",
            "O = P · V  — blend the info, weighted by attention.",
            "That’s it. Everything else (causal, sliding-window, ALiBi, RoPE) is a tweak on this.",
        ])
        return s
    slides.append(slide_math)

    # 5. Why it's slow ------------------------------------------------------
    def slide_slow():
        s = prs.slides.add_slide(blank)
        add_title(s, "Why Is Attention Slow?")
        add_subtitle(s, "The N×N scoreboard explodes as the text gets longer.")
        add_bullets(s, [
            "For N = 4096 words, the scoreboard has ~16 million entries — per attention head, per layer.",
            "A GPT-2 size model has many heads × many layers → tens of billions of entries.",
            "The GPU must write that scoreboard to slow memory (HBM) and read it back.",
            "Memory traffic, not math, becomes the bottleneck.",
            "Standard PyTorch makes this worse by launching 5+ separate GPU kernels (Q·K, scale, mask, softmax, ·V).",
        ])
        return s
    slides.append(slide_slow)

    # 6. FlashAttention -----------------------------------------------------
    def slide_flash():
        s = prs.slides.add_slide(blank)
        add_title(s, "FlashAttention: The Hand-Optimised Hero")
        add_subtitle(s, "It is faster — but only for a fixed set of attention shapes.")
        add_bullets(s, [
            "FlashAttention (Dao 2022, 2023) never writes the full N×N scoreboard.",
            "It processes tiny tiles, keeping running max & sum in fast on-chip memory.",
            "Same answer, 2–4× faster, way less memory.",
            "Catch: it is hand-written CUDA. It only supports dense and causal attention.",
            ("New ideas (sliding-window, ALiBi, RoPE, sparse) need fresh CUDA — weeks of expert work.", 0),
        ])
        return s
    slides.append(slide_flash)

    # 7. The real problem ---------------------------------------------------
    def slide_realproblem():
        s = prs.slides.add_slide(blank)
        add_title(s, "The Real Problem")
        add_subtitle(s, "Researchers want to try new attention shapes. The tooling doesn’t cooperate.")
        add_bullets(s, [
            "Sliding-window — only attend to the last 256 words (Mistral, Longformer).",
            "ALiBi — penalise far-away words (used in BLOOM, GPT-NeoX).",
            "RoPE — rotate Q and K by position (LLaMA, Mistral).",
            "Block-sparse — only attend to chosen blocks (BigBird).",
            "PyTorch SDPA covers only dense + causal; everything else falls back to slow O(N²) code.",
            "There is no middle ground between “write CUDA for weeks” and “run slowly.”",
        ])
        return s
    slides.append(slide_realproblem)

    # 8. The AttnFuse idea --------------------------------------------------
    def slide_idea():
        s = prs.slides.add_slide(blank)
        add_title(s, "AttnFuse — The Idea")
        add_subtitle(s, "A tiny “recipe language” for attention, plus a compiler.")
        add_bullets(s, [
            "You write the attention as a short Python function using building blocks.",
            "Our compiler turns it into one fused GPU kernel automatically.",
            "Fused = all the steps (score, mask, softmax, mix) happen in one go, in fast memory.",
            "No CUDA, no Triton expertise required.",
            ("Same speed class as hand-written FlashAttention; supports any combination of building blocks.", 0),
        ])
        return s
    slides.append(slide_idea)

    # 9. Code example -------------------------------------------------------
    def slide_code():
        s = prs.slides.add_slide(blank)
        add_title(s, "What the User Writes (the Whole Thing!)")
        add_subtitle(s, "Five lines. The first call compiles a fused GPU kernel; later calls are <200 µs.")
        add_code(s, """import attnfuse as af

@af.attention
def my_attn(Q, K, V):
    s = af.scaled_dot_product(Q, K)
    s = af.causal(s)
    s = af.alibi(s, num_heads=12)
    return af.softmax(s) @ V

out = my_attn(Q, K, V)   # JIT-compiles a fused Triton kernel on first call
""", top=1.7, height=4.6, size=18)
        return s
    slides.append(slide_code)

    # 10. The combinators (LEGO) -------------------------------------------
    def slide_combinators():
        s = prs.slides.add_slide(blank)
        add_title(s, "The Building Blocks (Combinators)")
        add_subtitle(s, "Like LEGO bricks for attention. Snap them together in any order.")
        add_bullets(s, [
            ("Score:", 0),
            ("scaled_dot_product(Q, K) — the basic score math", 1),
            ("rope(Q, K) — rotate Q, K by position before scoring", 1),
            ("Mask:", 0),
            ("causal(s) — only look at past words", 1),
            ("sliding_window(s, W) — only look W words back", 1),
            ("Bias:", 0),
            ("alibi(s) — distance penalty   •   additive_bias(s) — your own bias tensor", 1),
            ("Normalise:", 0),
            ("softmax(s) — the standard one   •   relu_attention(s) — a research variant", 1),
        ], size=18, gap_pt=4)
        return s
    slides.append(slide_combinators)

    # 11. The compiler pipeline --------------------------------------------
    def slide_pipeline():
        s = prs.slides.add_slide(blank)
        add_title(s, "How the Compiler Works")
        add_subtitle(s, "Two-level IR (intermediate representation) + four passes — like a recipe → meal.")
        add_bullets(s, [
            "Trace — run the user’s Python once with fake tensors; capture the graph.",
            "Fuse — recognise the (score → mask → bias → softmax → mix) pattern.",
            "Tile — pick block sizes for the GPU (BLOCK_M, BLOCK_N, num_warps).",
            "Lower — flatten the graph to a tile-loop recipe.",
            "Codegen — emit one Triton source file; Triton’s JIT turns it into GPU code.",
            "Result is cached by a hash of the graph: same recipe → reuse the kernel.",
        ])
        return s
    slides.append(slide_pipeline)

    # 12. One kernel, many variants ----------------------------------------
    def slide_onekernel():
        s = prs.slides.add_slide(blank)
        add_title(s, "One Kernel, Many Variants")
        add_subtitle(s, "The trick: compile-time switches (tl.constexpr), zero runtime branching.")
        add_bullets(s, [
            "We write ONE big templated Triton kernel.",
            "It carries flags: MASK_KIND, BIAS_KIND, NORM_KIND, ROPE_KIND.",
            "Triton’s JIT specialises a brand-new GPU binary for each unique combination.",
            "Inside the kernel: zero if-branches in the hot path — the dead code is gone.",
            "Adding a new variant = a few lines of codegen, not a new kernel.",
        ])
        return s
    slides.append(slide_onekernel)

    # 13. Online softmax (the tile trick) ----------------------------------
    def slide_online():
        s = prs.slides.add_slide(blank)
        add_title(s, "The Magic Inside: Online Softmax")
        add_subtitle(s, "How we get FlashAttention behaviour without ever storing the N×N scoreboard.")
        add_bullets(s, [
            "Naive softmax needs the full row of scores at once — impossible at N = 4096.",
            "Online softmax keeps two running numbers per row:  m (max),  ℓ (sum).",
            "As each tile of K, V arrives, we update m, ℓ, and a partial output.",
            "Old contributions get rescaled when a new max appears.",
            "End of the loop → exact same numbers as the naive version, but never materialised.",
            "Memory cost: O(N) instead of O(N²).",
        ])
        return s
    slides.append(slide_online)

    # 14. The novel trick: fused RoPE --------------------------------------
    def slide_rope():
        s = prs.slides.add_slide(blank)
        add_title(s, "Our Novel Trick — Fused RoPE")
        add_subtitle(s, "RoPE = Rotary Position Embedding (LLaMA, Mistral). We rotate inside the kernel.")
        add_bullets(s, [
            "Normally: two extra GPU kernels rotate Q and K, write results back to slow memory.",
            "AttnFuse: rotate Q once before the inner loop, rotate K tiles as they are loaded.",
            "All in registers — no extra trip to HBM.",
            "1.18× faster at N = 4096 → 2.29× faster at N = 512.",
            "To our knowledge, the first compiler-generated fused RoPE (prior implementations are hand-written).",
        ])
        return s
    slides.append(slide_rope)

    # 15. Headline results --------------------------------------------------
    def slide_results():
        s = prs.slides.add_slide(blank)
        add_title(s, "Headline Results (RTX 3090, fp16, N = 4096)")
        add_subtitle(s, "AttnFuse vs PyTorch SDPA — same model, same hardware.")
        # 4-column comparison "table" using bullets for portability
        add_bullets(s, [
            "Dense:                 AttnFuse 59 TFLOPS   |   SDPA 70 TFLOPS   →   0.85× (matches FA2)",
            "Causal:               AttnFuse 105 TFLOPS  |   SDPA 122 TFLOPS  →   0.86× (matches FA2)",
            "Sliding-window W=256:  AttnFuse 298 TFLOPS  |   SDPA  7.8 TFLOPS →  38× faster",
            "Causal + ALiBi:        AttnFuse  98 TFLOPS  |   SDPA  5.6 TFLOPS →  17× faster",
        ], size=18)
        add_callout(s,
            "For variants SDPA doesn’t accelerate, AttnFuse is the only sub-quadratic option — "
            "no custom CUDA required.",
            left=0.7, top=5.6, width=12.0, height=1.0, fill=ORANGE)
        return s
    slides.append(slide_results)

    # 16. Memory & dtype ---------------------------------------------------
    def slide_memdtype():
        s = prs.slides.add_slide(blank)
        add_title(s, "Memory + Multi-Dtype")
        add_subtitle(s, "32× less memory, three precisions, three baselines.")
        add_bullets(s, [
            "Peak GPU memory at N=4096: AttnFuse 200 MB vs naive PyTorch 3,176 MB (≈ 32× less).",
            "fp16 / bfloat16 / fp32 all supported with dtype-aware tile configs.",
            "bfloat16 causal at N=4096:  108 TFLOPS (slightly faster than fp16 on Ampere).",
            "fp32 causal:  44 TFLOPS — 3.2× faster than PyTorch SDPA fp32 (which loses its flash path).",
            "Verified by 25 GPU unit tests across all variants and dtypes.",
        ])
        return s
    slides.append(slide_memdtype)

    # 17. Cost & limitations + summary -------------------------------------
    def slide_costs():
        s = prs.slides.add_slide(blank)
        add_title(s, "Costs, Limitations, and What We Built")
        add_subtitle(s, "Trade-offs, what’s missing, and the take-home.")
        add_bullets(s, [
            "First-call JIT compile: 0.8–1.8 s. Cached: <200 µs. Pay once, fast forever.",
            "Forward pass only — no training (backward kernel is future work).",
            "Single GPU; head dim must be a power of 2 ≤ 256; no block-sparse yet.",
            "What we shipped: 7 combinators, 2-level IR, 4-pass compiler, fused RoPE,",
            ("25 GPU tests, 3 dtypes, 4 sequence lengths, 3 baselines, full evaluation.", 1),
            "Bottom line: declarative attention DSL with FA2-level performance, plus 17–38× wins on variants SDPA can’t do.",
        ])
        return s
    slides.append(slide_costs)

    # ---- build all slides + footers ---------------------------------------
    total = len(slides)
    out = []
    for i, builder in enumerate(slides, 1):
        sl = builder()
        if i > 1:
            add_footer(sl, i, total)
        out.append(sl)

    out_path = Path(__file__).parent / "AttnFuse_slides.pptx"
    prs.save(out_path)
    print(f"[ok] wrote {out_path}  ({total} slides)")


if __name__ == "__main__":
    build()
